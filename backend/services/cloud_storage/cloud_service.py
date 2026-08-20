# This part of the code combines mock_cloud_service and s3_service to provide a unified interface for cloud operations. It allows the application to switch between mock (GCP and Azure) and real S3 services
import io
import re
import zipfile
from pathlib import Path

from backend.services.cloud_storage.s3_service import (
    list_s3_objects,
    read_s3_file,
    write_s3_file,
    get_s3_object_metadata,
)

from backend.services.cloud_storage.mock_cloud_service import (
    list_mock_objects,
    read_mock_file,
    write_mock_file,
    get_mock_object_metadata,
)

# ---------------------------------------------------------------------------
# Optional library imports — graceful degradation if not installed
# ---------------------------------------------------------------------------
try:
    import pdfplumber as _pdfplumber
    _HAS_PDFPLUMBER = True
except ImportError:
    _HAS_PDFPLUMBER = False

try:
    import PyPDF2 as _PyPDF2
    _HAS_PYPDF2 = True
except ImportError:
    _HAS_PYPDF2 = False

try:
    from docx import Document as _DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    import openpyxl as _openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

try:
    import xlrd as _xlrd
    _HAS_XLRD = True
except ImportError:
    _HAS_XLRD = False

try:
    from pptx import Presentation as _Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

_MAX_CHARS = 500_000


def _printable_strings(data: bytes) -> str:
    """Last-resort extractor: pull printable ASCII runs out of any binary blob."""
    out: list[str] = []
    run: list[int] = []
    for byte in data[: _MAX_CHARS * 4]:
        if byte == 9 or byte == 10 or byte == 13 or 32 <= byte <= 126:
            run.append(byte)
        else:
            if len(run) >= 4:
                out.append(bytes(run).decode("ascii", errors="ignore"))
            run = []
    if len(run) >= 4:
        out.append(bytes(run).decode("ascii", errors="ignore"))
    return "\n".join(out)[:_MAX_CHARS]


def _strip_xml(xml_bytes: bytes) -> str:
    text = re.sub(rb"<[^>]+>", b" ", xml_bytes).decode("utf-8", errors="ignore")
    return " ".join(text.split())


def _read_zip_xml(data: bytes) -> str:
    """Extract text from Office Open XML / ODF ZIP containers."""
    try:
        parts: list[str] = []
        written = 0
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in zf.namelist():
                if not name.lower().endswith(".xml"):
                    continue
                try:
                    plain = _strip_xml(zf.read(name))
                except Exception:
                    continue
                if not plain:
                    continue
                chunk = f"\n[{name}]\n{plain}\n"
                if written + len(chunk) > _MAX_CHARS:
                    break
                parts.append(chunk)
                written += len(chunk)
        return "\n".join(parts).strip()
    except Exception:
        return ""


def extract_text(data: bytes, filename: str = "") -> str:
    """
    Extract plain text from raw file bytes so PII detection works on
    the content regardless of file format.

    Supported formats
    -----------------
    Plain text : .txt .csv .tsv .log .md .json .xml .yaml .html .sql …
    PDF        : .pdf  — pdfplumber → PyPDF2 → printable-string fallback
    Word       : .docx — python-docx; .doc — printable strings
    Excel      : .xlsx/.xlsm — openpyxl; .xls — xlrd; .csv — plain text
    PowerPoint : .pptx — python-pptx; .ppt — printable strings
    OpenDoc    : .odt .ods .odp — ZIP+XML strip
    Rich Text  : .rtf — strip control words
    E-mail     : .eml — stdlib email parser
    """
    if not data:
        return ""

    ext = Path(filename).suffix.lower() if filename else ""

    # ---- PDF ----
    if ext == ".pdf":
        if _HAS_PDFPLUMBER:
            try:
                pages: list[str] = []
                with _pdfplumber.open(io.BytesIO(data)) as pdf:
                    for page in pdf.pages:
                        pages.append(page.extract_text() or "")
                        if sum(len(p) for p in pages) >= _MAX_CHARS:
                            break
                result = "\n".join(pages).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        if _HAS_PYPDF2:
            try:
                reader = _PyPDF2.PdfReader(io.BytesIO(data))
                pages = [p.extract_text() or "" for p in reader.pages]
                result = "\n".join(pages).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        return _printable_strings(data)

    # ---- Word DOCX ----
    if ext == ".docx":
        if _HAS_DOCX:
            try:
                doc = _DocxDocument(io.BytesIO(data))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text)
                result = "\n".join(parts).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        return _read_zip_xml(data) or _printable_strings(data)

    # ---- Excel XLSX / XLSM ----
    if ext in (".xlsx", ".xlsm"):
        if _HAS_OPENPYXL:
            try:
                wb = _openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            parts.append(", ".join(cells))
                        if sum(len(p) for p in parts) >= _MAX_CHARS:
                            break
                result = "\n".join(parts).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        return _read_zip_xml(data) or _printable_strings(data)

    # ---- Excel XLS (legacy binary) ----
    if ext == ".xls":
        if _HAS_XLRD:
            try:
                wb = _xlrd.open_workbook(file_contents=data)
                parts = []
                for sheet in wb.sheets():
                    for row_idx in range(sheet.nrows):
                        cells = [
                            str(sheet.cell_value(row_idx, col)).strip()
                            for col in range(sheet.ncols)
                            if str(sheet.cell_value(row_idx, col)).strip()
                        ]
                        if cells:
                            parts.append(", ".join(cells))
                result = "\n".join(parts).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        return _printable_strings(data)

    # ---- PowerPoint PPTX ----
    if ext == ".pptx":
        if _HAS_PPTX:
            try:
                prs = _Presentation(io.BytesIO(data))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                result = "\n".join(parts).strip()
                if result:
                    return result[:_MAX_CHARS]
            except Exception:
                pass
        return _read_zip_xml(data) or _printable_strings(data)

    # ---- Legacy binary Office (.doc, .ppt) ----
    if ext in (".doc", ".ppt"):
        return _printable_strings(data)

    # ---- OpenDocument (.odt, .ods, .odp) ----
    if ext in (".odt", ".ods", ".odp"):
        return _read_zip_xml(data) or _printable_strings(data)

    # ---- RTF ----
    if ext == ".rtf":
        try:
            text = data.decode("latin-1", errors="ignore")
            text = re.sub(r"\{[^{}]*\}", "", text)
            text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", text)
            text = text.replace("{", " ").replace("}", " ")
            return " ".join(text.split())[:_MAX_CHARS]
        except Exception:
            return _printable_strings(data)

    # ---- E-mail ----
    if ext == ".eml":
        try:
            from email import message_from_bytes as _mfb
            msg = _mfb(data)
            parts = []
            for header in ("From", "To", "Cc", "Subject", "Date"):
                val = msg.get(header, "")
                if val:
                    parts.append(f"{header}: {val}")
            if msg.is_multipart():
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct in ("text/plain", "text/html"):
                        payload = part.get_payload(decode=True)
                        if isinstance(payload, bytes):
                            charset = part.get_content_charset() or "utf-8"
                            chunk = payload.decode(charset, errors="ignore")
                            if ct == "text/html":
                                chunk = re.sub(r"<[^>]+>", " ", chunk)
                            parts.append(chunk)
            else:
                payload = msg.get_payload(decode=True)
                if isinstance(payload, bytes):
                    parts.append(payload.decode(msg.get_content_charset() or "utf-8", errors="ignore"))
            return "\n".join(parts)[:_MAX_CHARS]
        except Exception:
            return _printable_strings(data)

    # ---- Plain text (default) ----
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return data.decode(encoding)[:_MAX_CHARS]
        except (UnicodeDecodeError, ValueError):
            continue
    return _printable_strings(data)


def list_cloud_objects(org_id: str | None = None):

    aws_objects = list_s3_objects()

    mock_objects = list_mock_objects(org_id=org_id)

    return aws_objects + mock_objects


def list_files(org_id: str | None = None):

    return [
        obj["file"]
        for obj in list_cloud_objects(org_id=org_id)
    ]


def read_file(path):
    # REAL AWS S3 — returns bytes, run through extractor
    if isinstance(path, str) and path.startswith("s3://"):
        raw = read_s3_file(path)
        if isinstance(raw, bytes):
            return extract_text(raw, filename=path)
        return raw  # already a string (legacy caller)

    # MOCK GCP / AZURE — local file path, read bytes then extract
    try:
        data = Path(path).read_bytes()
        return extract_text(data, filename=str(path))
    except Exception:
        return read_mock_file(path)  # fallback to original plain-text read


def write_file(path, content):

    # REAL AWS S3
    if isinstance(path, str) and path.startswith("s3://"):

        return write_s3_file(
            path,
            content
        )

    # MOCK GCP / AZURE
    return write_mock_file(
        path,
        content
    )


def get_object_metadata(path):

    # REAL AWS S3
    if isinstance(path, str) and path.startswith("s3://"):

        return get_s3_object_metadata(path)

    # MOCK GCP / AZURE
    return get_mock_object_metadata(path)
